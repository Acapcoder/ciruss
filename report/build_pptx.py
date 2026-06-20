"""Build a 6-slide CIRRUS presentation (red and white theme, no em dashes)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(__file__)
ASSETS = os.path.join(HERE, "assets")
OUT = os.path.join(HERE, "CIRRUS_Presentation.pptx")

RED = RGBColor(0xBE, 0x12, 0x3C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x11, 0x11, 0x11)
GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF7, 0xF2, 0xF4)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
    return s


def textbox(s, left, top, width, height):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tb.text_frame.word_wrap = True
    return tb.text_frame


def set_run(r, text, size, color=BLACK, bold=False, italic=False):
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = "Calibri"


def title_band(s, title):
    band = s.shapes.add_shape(1, 0, 0, SW, Inches(1.15))
    band.fill.solid(); band.fill.fore_color.rgb = RED
    band.line.fill.background(); band.shadow.inherit = False
    tf = band.text_frame
    tf.margin_left = Inches(0.6); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    set_run(p.add_run(), title, 30, WHITE, bold=True)


def bullets(s, items, left=0.7, top=1.5, width=7.6, height=5.5, size=19):
    tf = textbox(s, left, top, width, height)
    for i, (text, lvl) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        p.space_after = Pt(10)
        prefix = "" if lvl == 0 else ""
        set_run(p.add_run(), ("•  " if lvl == 0 else "–  ") + text, size if lvl == 0 else size - 3,
                BLACK if lvl == 0 else GRAY, bold=(lvl == 0))


def png_size(path):
    import struct
    with open(path, "rb") as f:
        head = f.read(24)
    return struct.unpack(">II", head[16:24])  # width, height from IHDR


def picture_fit(s, path, left, top, max_w, max_h):
    try:
        iw, ih = png_size(path)
    except Exception:
        iw, ih = (1360, 900)
    ratio = min(max_w / iw, max_h / ih)
    w = Emu(int(iw * ratio * 914400))
    h = Emu(int(ih * ratio * 914400))
    pic = s.shapes.add_picture(path, Inches(left), Inches(top), width=w, height=h)
    pic.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD); pic.line.width = Pt(0.75)
    return pic


# ---------------- Slide 1: Title ----------------
s = slide()
bar = s.shapes.add_shape(1, 0, Inches(2.55), SW, Inches(0.08))
bar.fill.solid(); bar.fill.fore_color.rgb = RED; bar.line.fill.background(); bar.shadow.inherit = False
tf = textbox(s, 1.0, 2.7, 11.3, 1.2)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
set_run(p.add_run(), "CIRRUS", 60, RED, bold=True)
tf2 = textbox(s, 1.0, 3.9, 11.3, 1.0)
p = tf2.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
set_run(p.add_run(), "Smart Multi-Cloud Storage Pooling and Compression Platform", 22, BLACK)
tf3 = textbox(s, 1.0, 5.2, 11.3, 1.2)
for i, line in enumerate(["Project Presentation", "Team Members: [Member 1], [Member 2], [Member 3], [Member 4]"]):
    p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
    p.alignment = PP_ALIGN.CENTER
    set_run(p.add_run(), line, 16, GRAY, bold=(i == 0))

# ---------------- Slide 2: Problem and Solution ----------------
s = slide(); title_band(s, "The Problem and Our Solution")
bullets(s, [
    ("The Problem", 0),
    ("Free cloud storage is small and scattered, for example 15 GB on Drive and 2 GB on Dropbox.", 1),
    ("Users hold several accounts but cannot use the free space as one pool.", 1),
    ("No simple consumer tool unifies these free allocations.", 1),
    ("Our Solution: CIRRUS", 0),
    ("Connect many free accounts and treat them as one combined storage box.", 1),
    ("Compress every file, then auto-route it to the cloud with the most free space.", 1),
    ("Organize everything in one virtual folder tree, regardless of which cloud holds it.", 1),
], width=11.9)

# ---------------- Slide 3: Architecture and Tech Stack ----------------
s = slide(); title_band(s, "System Architecture")
bullets(s, [
    ("Three tier design", 0),
    ("Frontend: React and Vite, hosted on Vercel.", 1),
    ("Backend: Python FastAPI, hosted on Amazon Web Services.", 1),
    ("Database: managed PostgreSQL on Supabase (metadata and encrypted tokens only).", 1),
    ("Storage: files live on the user's Google Drive and Dropbox.", 1),
    ("Data flow", 0),
    ("Browser sends file to backend, backend compresses it, picks the emptiest cloud, uploads, and records metadata.", 1),
], width=11.9)

# ---------------- Slide 4: Key Features (with dashboard image) ----------------
s = slide(); title_band(s, "Key Features")
bullets(s, [
    ("Storage pooling: combined free space shown as one box.", 0),
    ("Compression and smart routing of every upload.", 0),
    ("Virtual folders and bulk multi-file upload.", 0),
    ("Secure sign in with encrypted cloud credentials.", 0),
    ("Live dashboard with usage, savings, and activity log.", 0),
], left=0.6, top=1.5, width=6.0, size=18)
picture_fit(s, os.path.join(ASSETS, "03_dashboard.png"), 6.9, 1.5, 6.1, 5.4)

# ---------------- Slide 5: User Interface ----------------
s = slide(); title_band(s, "The Platform in Action")
picture_fit(s, os.path.join(ASSETS, "04_connections.png"), 0.5, 1.45, 6.1, 5.6)
picture_fit(s, os.path.join(ASSETS, "05_files.png"), 6.85, 1.45, 6.1, 5.6)
cap = textbox(s, 0.5, 6.95, 12.4, 0.4)
p = cap.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
set_run(p.add_run(), "Connecting clouds (left) and the folder based file manager with bulk upload (right).", 12, GRAY, italic=True)

# ---------------- Slide 6: Team and Conclusion ----------------
s = slide(); title_band(s, "Team Contributions and Conclusion")
bullets(s, [
    ("Member 1: Backend core, database schema, compression and routing engine.", 0),
    ("Member 2: Authentication, password hashing, credential encryption, audit logs.", 0),
    ("Member 3: Google Drive and Dropbox integration and token handling.", 0),
    ("Member 4: React interface, red and white design, and full deployment.", 0),
    ("Conclusion", 0),
    ("CIRRUS turns scattered free storage into one secure, compressed, and easy to use box.", 1),
], width=11.9, size=18)

prs.save(OUT)
print("Presentation written to", OUT, "with", len(prs.slides._sldIdLst), "slides,",
      round(os.path.getsize(OUT) / 1024), "KB")
